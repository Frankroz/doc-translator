import os
import threading
import time
import tkinter as tk
from time import sleep
import copy
import openpyxl
from openpyxl.styles import Alignment
from openpyxl.cell.cell import MergedCell
from docx import Document
from pptx import Presentation
from deep_translator import GoogleTranslator
import customtkinter as ctk
from tkinter import filedialog, messagebox
import random
import gc
import subprocess
import requests
import json
from concurrent.futures import ThreadPoolExecutor, as_completed
import re

try:
    from win32com.client import Dispatch
    import win32com.client as win32
except ImportError:
    win32 = None

ctk.set_appearance_mode("Dark")
ctk.set_default_color_theme("blue")

class DocTranslatorPro(ctk.CTk):
    def __init__(self):
        super().__init__()

        self.title("With Love ❤︎ - Doc Translator")

        self.after(0, lambda: self.state('zoomed'))

        self.geometry("1000x820")
        self.URL="https://lovenotesgroup-cjfwccfqdkggh7g4.canadacentral-01.azurewebsites.net"

        self.languages = {
            "Español": "es",
            "Inglés": "en",
            "Chino": "zh-CN",
            "Auto-detectar": "auto"
        }
        
        self.is_cancelled = False
        self.files_path_set = set()
        self.files_to_process = []
        self.keep_original = tk.BooleanVar(value=True)
        
        self.total_units = 0
        self.done_units = 0
        self.start_time = 0

        self._build_ui()
        self._animate_heart()

    def _build_ui(self):
        # --- CONFIGURACIÓN DE ESCALADO GLOBAL ---
        # Columna principal se expande
        self.grid_columnconfigure(0, weight=1)
        
        # Filas: La fila 3 (scrollable frame) se lleva el espacio sobrante
        self.grid_rowconfigure(0, weight=0) # Cabecera
        self.grid_rowconfigure(1, weight=0) # Idiomas
        self.grid_rowconfigure(2, weight=0) # Botones carga
        self.grid_rowconfigure(3, weight=1) # LISTA (EXPANDIBLE)
        self.grid_rowconfigure(4, weight=0) # Panel inferior

        # 1. CABECERA
        self.head = ctk.CTkLabel(self, text="TRANSLATOR FOR THE BEST TRANSLATOR", 
                                font=ctk.CTkFont(size=32, weight="bold"))
        self.head.grid(row=0, column=0, pady=(40, 20), sticky="ew")

        # 2. SECCIÓN DE IDIOMAS
        self.lang_frame = ctk.CTkFrame(self, fg_color="transparent")
        self.lang_frame.grid(row=1, column=0, pady=10)
        # No necesita sticky="nsew" para mantenerse agrupado al centro

        ctk.CTkLabel(self.lang_frame, text="De:").grid(row=0, column=0, padx=5)
        self.combo_from = ctk.CTkOptionMenu(self.lang_frame, values=list(self.languages.keys()))
        self.combo_from.set("Auto-detectar")
        self.combo_from.grid(row=0, column=1, padx=10)

        ctk.CTkLabel(self.lang_frame, text="A:").grid(row=0, column=2, padx=5)
        self.combo_to = ctk.CTkOptionMenu(self.lang_frame, values=["Español", "Inglés", "Chino"])
        self.combo_to.set("Español")
        self.combo_to.grid(row=0, column=3, padx=10)    

        # 3. BOTONES DE CARGA Y OPCIONES
        self.up_frame = ctk.CTkFrame(self, fg_color="transparent")
        self.up_frame.grid(row=2, column=0, padx=20, pady=10)
        # Configuramos columnas del up_frame para centrar contenido
        self.up_frame.grid_columnconfigure((0, 1, 2), weight=1)

        self.btn_fold = ctk.CTkButton(self.up_frame, text="📂 Carpeta", command=self.add_folder, fg_color="#27ae60", width=150)
        self.btn_fold.grid(row=0, column=0, padx=10, pady=5)

        self.btn_one = ctk.CTkButton(self.up_frame, text="📄 Archivo", command=self.add_file, fg_color="#2980b9", width=150)
        self.btn_one.grid(row=0, column=1, padx=10, pady=5)

        self.btn_clr = ctk.CTkButton(self.up_frame, text="🗑️ Limpiar", command=self.clear, fg_color="#7f8c8d", width=150)
        self.btn_clr.grid(row=0, column=2, padx=10, pady=5)

        # FIX: Cambiado pack() por grid() para evitar conflictos
        self.check_keep = ctk.CTkCheckBox(self.up_frame, text="Mantener Texto Original", 
                                          variable=self.keep_original,
                                          font=("Segoe UI", 12))
        self.check_keep.grid(row=1, column=0, columnspan=3, pady=15) # Centrado debajo de los botones

        # 4. LISTA DE ARCHIVOS
        self.list_frame = ctk.CTkScrollableFrame(self, label_text="Cola de Documentos (Word, Excel, PPT)", 
                                                 label_font=ctk.CTkFont(weight="bold"))
        self.list_frame.grid(row=3, column=0, padx=50, pady=20, sticky="nsew")

        # 5. PANEL INFERIOR (CONTROLES Y PROGRESO)
        self.bot_panel = ctk.CTkFrame(self, corner_radius=20, fg_color="#1a1a1a")
        self.bot_panel.grid(row=4, column=0, padx=40, pady=30, sticky="ew") # Se expande a lo ancho
        self.bot_panel.grid_columnconfigure(0, weight=1) # Centra contenido interno

        self.lbl_status = ctk.CTkLabel(self.bot_panel, text="Esperando archivos...", font=ctk.CTkFont(size=16))
        self.lbl_status.grid(row=0, column=0, pady=(20, 5), sticky="ew")

        self.lbl_msg = ctk.CTkLabel(self.bot_panel, text="", font=ctk.CTkFont(size=16))
        self.lbl_msg.grid(row=1, column=0, pady=(0, 5), sticky="ew")

        self.lbl_eta = ctk.CTkLabel(self.bot_panel, text="ETA: --:--", text_color="#95a5a6")
        self.lbl_eta.grid(row=2, column=0, pady=(0, 10), sticky="ew")

        # Barra de progreso responsiva: eliminamos width fijo de 800 y usamos sticky
        self.bar = ctk.CTkProgressBar(self.bot_panel, height=12) 
        self.bar.set(0)
        self.bar.grid(row=3, column=0, pady=10, padx=50, sticky="ew") # Se ajusta al ancho del panel

        # BOTONES DE CONTROL FINALES
        self.ctrl_btns = ctk.CTkFrame(self.bot_panel, fg_color="transparent")
        self.ctrl_btns.grid(row=4, column=0, pady=(10, 25))

        self.btn_go = ctk.CTkButton(self.ctrl_btns, text="🚀 INICIAR TRADUCCIÓN GLOBAL", command=self.start, 
                                    width=300, height=50, font=ctk.CTkFont(size=16, weight="bold"), state="disabled")
        self.btn_go.grid(row=0, column=0, padx=10)

        self.btn_stop = ctk.CTkButton(self.ctrl_btns, text="🛑 CANCELAR", command=self.stop, 
                                    width=150, height=50, fg_color="#c0392b", state="disabled")
        self.btn_stop.grid(row=0, column=1, padx=10)
    
    def _animate_heart(self):
        beats = ["❤️", "💓", "❤️", "💓"] 
        delays = [0.2, 0.4, 0.2, 0.8]
        def run():
            idx = 0
            while True:
                try:
                    h = beats[idx % len(beats)]
                    self.title(f"Hecho con amor para ti {h} - Doc Translator")
                    time.sleep(delays[idx % len(delays)])
                    idx += 1
                except: break
        threading.Thread(target=run, daemon=True).start()

    def add_folder(self):
        f = filedialog.askopenfilename(
            title="Navega hasta la carpeta y selecciona cualquier archivo dentro para usar esa carpeta",
            filetypes=[("Archivos Office", "*.docx *.doc *.xlsx *.xls *.pptx *.ppt")]
        )
        
        if f:
            # Extraemos la ruta de la carpeta del archivo seleccionado
            folder_path = os.path.dirname(f)
            
            exts = ('.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt')
            files_found = False
            
            for file in os.listdir(folder_path):
                if file.lower().endswith(exts) and "_BILINGUAL_" not in file.upper() and not file.startswith('~$'):
                    self.register_file(os.path.join(folder_path, file))
                    files_found = True
            
            if not files_found:
                messagebox.showwarning("Atención", "No se encontraron archivos compatibles en esta carpeta.")

    def add_file(self):
        f = filedialog.askopenfilename(filetypes=[("Archivos Office", "*.docx *.doc *.xlsx *.xls *.pptx *.ppt")])
        if f and "_BILINGUAL_" not in os.path.basename(f).upper():
            self.register_file(f)

    def register_file(self, path):
        abs_p = os.path.abspath(path)
        if abs_p not in self.files_path_set:
            self.files_path_set.add(abs_p)
            
            row_frame = ctk.CTkFrame(self.list_frame, fg_color="transparent")
            row_frame.pack(fill="x", padx=5, pady=2)
            
            v = ctk.BooleanVar(value=True)
            chk = ctk.CTkCheckBox(row_frame, text=os.path.basename(path), variable=v)
            chk.pack(side="left", padx=10, fill="x", expand=True)

            btn_del = ctk.CTkButton(row_frame, text="✕", width=30, 
                                   fg_color="#c0392b", hover_color="#e74c3c",
                                   command=lambda p=abs_p, f=row_frame: self.remove_file_from_queue(p, f))
            btn_del.pack(side="right", padx=2)
            
            btn_up = ctk.CTkButton(row_frame, text="▲", width=30, fg_color="#34495e", 
                                   command=lambda: self.move_file(abs_p, -1))
            btn_up.pack(side="right", padx=2)
            
            btn_down = ctk.CTkButton(row_frame, text="▼", width=30, fg_color="#34495e", 
                                     command=lambda: self.move_file(abs_p, 1))
            btn_down.pack(side="right", padx=2)
            
            self.files_to_process.append({
                "path": abs_p, 
                "var": v, 
                "frame": row_frame
            })
            self.btn_go.configure(state="normal")

    def clear(self):
        for w in self.list_frame.winfo_children(): w.destroy()
        self.files_to_process, self.files_path_set = [], set()
        self.btn_go.configure(state="disabled")

    def stop(self):
        self.is_cancelled = True
        self.lbl_status.configure(text="Cancelando...", text_color="#e74c3c")
        self.btn_go.configure(state="normal")
        self.btn_stop.configure(state="disabled")
        self.check_keep(state="enabled")

    def start(self):
        self.current_src_lang = self.languages[self.combo_from.get()]
        self.current_tgt_lang = self.languages[self.combo_to.get()]
        lang_src_name = self.combo_from.get()
        lang_tgt_name = self.combo_to.get()

        self.start_time_deb = time.perf_counter()
        self.api_counter = 0

        # Mensaje de advertencia de "Peligro"
        msg = (f"⚠️ ATENCIÓN: Vas a traducir de {lang_src_name} a {lang_tgt_name}.\n\n"
           f"Los archivos se guardarán como '{lang_tgt_name}_BILINGUAL_...'\n"
           "¿Deseas comenzar?")
    
        if not messagebox.askyesno("Confirmar Idiomas", msg):
            return
        
        self.translation_cache = self.load_cache()

        print("INICIANDO PROCESO DE TRADUCCIÓN")
        sel = [item["path"] for item in self.files_to_process if item["var"].get()]
        if not sel: return

        self.translator = GoogleTranslator(source=self.current_src_lang, target=self.current_tgt_lang)
        self.is_cancelled = False
        self.btn_go.configure(state="disabled")
        self.btn_stop.configure(state="normal")
        #self.check_keep(state="disabled")

        self.selected_text = ""
        try:
            response = requests.get(f"{self.URL}/quote/ending", timeout=2)
            self.selected_text = response.json().get("text")
        except:
            self.selected_text = "✨ ¡Todo listo! ✨\n\nEspero que esto te ahorre\nmucho tiempo hoy.\n\nTe quiero muchoooo!."

        self.lbl_msg.configure(text=self.type_love_note())
        threading.Thread(target=self.main_loop, args=(sel,), daemon=True).start()

    def move_file(self, path, direction):
        idx = next(i for i, x in enumerate(self.files_to_process) if x["path"] == path)
        new_idx = idx + direction
        
        if 0 <= new_idx < len(self.files_to_process):
            self.files_to_process[idx], self.files_to_process[new_idx] = \
                self.files_to_process[new_idx], self.files_to_process[idx]
            
            self.refresh_file_list()

    def refresh_file_list(self):    
        for item in self.files_to_process:
            item["frame"].pack_forget()
        
        for item in self.files_to_process:
            item["frame"].pack(fill="x", padx=5, pady=2)

    def main_loop(self, files):
        self.lbl_status.configure(text="Analizando archivos y formatos legados...", text_color="#ffffff")
        self.total_units = 0
        self.done_units = 0
        ready_queue = []

        try:
            for p in files:
                if self.is_cancelled: break
                units, temp_p = self.prepare_document(p)
                self.total_units += units
                ready_queue.append((p, temp_p))

            if self.total_units == 0: self.total_units = 1
            self.start_time = time.time()

            for orig, temp in ready_queue:
                if self.is_cancelled: break
                self.lbl_status.configure(text=f"Procesando: {os.path.basename(orig)}", text_color="#ecf0f1")
                
                try:
                    if temp.endswith('.docx'):
                        self.translate_word(orig, temp)
                    elif temp.endswith('.xlsx'):
                        if win32:
                            self.lbl_status.configure(text=f"Traduciendo formas en: {os.path.basename(orig)}", text_color="#ecf0f1")
                            self.translate_shapes_via_win32(temp)
                        self.lbl_status.configure(text=f"Procesando: {os.path.basename(orig)}", text_color="#ecf0f1")
                        self.translate_excel(orig, temp)
                    elif temp.endswith('.pptx'):
                        if win32:
                            self.lbl_status.configure(text=f"Traduciendo formas en PPT: {os.path.basename(orig)}")
                            self.translate_pptx_shapes_via_win32(temp)
                        
                        self.translate_pptx(orig, temp)
                except Exception as e:
                    messagebox.showerror("Error", f"Error procesando {orig}: {e}")
                finally:
                    if "_TEMP_" in temp and os.path.exists(temp):
                        os.remove(temp)

        finally:
            self.bar.set(1 if not self.is_cancelled else 0)
            self.save_cache()
            self.lbl_status.configure(text="¡Completado!" if not self.is_cancelled else "Cancelado", text_color="#2ecc71" if not self.is_cancelled else "#e74c3c")
            self.btn_go.configure(state="normal"); self.btn_stop.configure(state="disabled")
            self.after(0, self.show_love_note)
            print("Finalizado en ", time.perf_counter()-self.start_time_deb, "segundos. Llamadas API:", self.api_counter)

    def prepare_document(self, path):
        work_path = path
        if path.lower().endswith(('.doc', '.xls', '.ppt')) and win32:
            import pythoncom
            pythoncom.CoInitialize()
            app = None
            try:
                abs_path = os.path.abspath(path)
                if path.lower().endswith('.doc'):
                    app = win32.Dispatch('Word.Application')
                    doc = app.Documents.Open(abs_path)
                    work_path = abs_path + "_TEMP_.docx"
                    doc.SaveAs2(work_path, FileFormat=16)
                    doc.Close()
                elif path.lower().endswith('.xls'):
                    app = win32.Dispatch('Excel.Application')
                    app.DisplayAlerts = False
                    wb = app.Workbooks.Open(abs_path)
                    work_path = abs_path + "_TEMP_.xlsx"
                    wb.SaveAs(work_path, FileFormat=51)
                    wb.Close(False)
                elif path.lower().endswith('.ppt'):
                    app = win32.Dispatch('PowerPoint.Application')
                    pres = app.Presentations.Open(abs_path, WithWindow=False)
                    work_path = abs_path + "_TEMP_.pptx"
                    pres.SaveAs(work_path, FileFormat=24)
                    pres.Close()
            except Exception as e:
                messagebox.showerror("Error", f"Error conversión legacy: {e}")
            finally:
                if app:
                    app.Quit()
                    app = None
                pythoncom.CoUninitialize()
                self.cleanup_janitor()
                time.sleep(2.0)

        cnt = 0
        try:
            if work_path.lower().endswith('.docx'):
                d = Document(work_path)
                cnt = len(d.paragraphs) + sum(len(r.cells) for t in d.tables for r in t.rows)
            elif work_path.lower().endswith('.xlsx'):
                for _ in range(3):
                    try:
                        wb_count = openpyxl.load_workbook(work_path, read_only=True)
                        for s in wb_count.worksheets: 
                            cnt += (s.max_row if s.max_row else 1)
                        wb_count.close()
                        break
                    except:
                        time.sleep(1)
            elif work_path.lower().endswith('.pptx'):
                prs = Presentation(work_path)
                cnt = len(prs.slides)
        except: cnt = 1
        return cnt, work_path

    def translate_pptx(self, orig, temp):
        prs = Presentation(temp)
        target_lang_name = self.combo_to.get().upper()
        out_name = f"{target_lang_name}_BILINGUAL_{os.path.basename(orig)}"
        output_path = os.path.join(os.path.dirname(orig), out_name)

        # 1. Recolectar todas las fuentes de texto (incluyendo grupos)
        objects_to_translate = []
        for slide in prs.slides:
            # FIX: Usamos el helper recursivo para no saltarnos grupos
            all_shapes = self.get_all_shapes(slide.shapes)
            
            for shape in all_shapes:
                if hasattr(shape, "text_frame") and self.is_valid_for_translation(shape.text):
                    objects_to_translate.append(shape)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if self.is_valid_for_translation(cell.text):
                                objects_to_translate.append(cell)

        # 2. Configuración de lotes
        batch_size = 15
        batches = [objects_to_translate[i : i + batch_size] for i in range(0, len(objects_to_translate), batch_size)]
        
        with ThreadPoolExecutor(max_workers=4) as executor:
            future_to_batch = {
                executor.submit(self.process_batch, [obj.text for obj in batch]): (i, batch) 
                for i, batch in enumerate(batches)
            }

            for future in as_completed(future_to_batch):
                if self.is_cancelled:
                    executor.shutdown(wait=False)
                    return

                batch_index, current_batch_objects = future_to_batch[future]
                try:
                    translated_results = future.result()
                    
                    for idx, obj in enumerate(current_batch_objects):
                        res = translated_results[idx] if idx < len(translated_results) else ""
                        
                        if hasattr(obj, "text_frame"):
                            if obj.text_frame.paragraphs:
                                last_para = obj.text_frame.paragraphs[-1]
                                
                                # Capturar estilo original del primer run disponible
                                source_run = last_para.runs[0] if last_para.runs else None
                                
                                if self.keep_original.get():
                                    new_run = last_para.add_run()
                                    new_run.text = f"\n{res}"
                                    # COPIA DE ESTILOS EXPLÍCITA
                                    if source_run:
                                        self.copy_font_style(source_run.font, new_run.font)
                                else:
                                    obj.text_frame.text = res
                                    # Aplicar estilo al texto reemplazado
                                    if source_run and obj.text_frame.paragraphs[0].runs:
                                        self.copy_font_style(source_run.font, obj.text_frame.paragraphs[0].runs[0].font)
                        
                        else: # Celda de Tabla
                            if self.keep_original.get():
                                obj.text = f"{obj.text}\n{res}"
                            else:
                                obj.text = res
                        
                        self.upd_prog()
                    
                    prs.save(output_path)
                    
                except Exception as e:
                    messagebox.showerror("Error", f"Error en lote {batch_index}: {e}")

    def copy_font_style(self, source_font, target_font):
        """Helper to duplicate visual attributes."""
        try:
            target_font.name = source_font.name
            target_font.size = source_font.size
            target_font.bold = source_font.bold
            target_font.italic = source_font.italic
            target_font.underline = source_font.underline
            if source_font.color and hasattr(source_font.color, 'rgb'):
                target_font.color.rgb = source_font.color.rgb
        except:
            pass

    def translate_word(self, orig, temp):
        doc = Document(temp)
        target_lang_name = self.combo_to.get().upper()
        out_name = f"{target_lang_name}_BILINGUAL_{os.path.basename(orig)}"
        if orig.lower().endswith('.doc'): out_name += "x"
        output_path = os.path.join(os.path.dirname(orig), out_name)

        # 1. Recolectar todos los elementos con texto
        elements_to_translate = []
        
        # Párrafos normales
        for p in doc.paragraphs:
            if self.is_valid_for_translation(p.text):
                elements_to_translate.append(p)
                
        # Texto en tablas
        for table in doc.tables:
            processed_cells = set()
            for row in table.rows:
                for cell in row.cells:
                    if cell._tc not in processed_cells and self.is_valid_for_translation(cell.text):
                        elements_to_translate.append(cell)
                        processed_cells.add(cell._tc)

        # 2. Procesar en lotes de 15-20
        batch_size = 15
        for i in range(0, len(elements_to_translate), batch_size):
            if self.is_cancelled: return
            
            current_batch = elements_to_translate[i : i + batch_size]
            texts_to_send = [el.text for el in current_batch]
            
            # Traducir el lote (usando tu caché interna)
            translated_texts = self.process_batch(texts_to_send)
            
            # 3. Aplicar traducciones de vuelta al documento
            for idx, el in enumerate(current_batch):
                res = translated_texts[idx] if idx < len(translated_texts) else ""
                
                if isinstance(el, type(doc.paragraphs[0])): # Es un párrafo
                    run = el.add_run(f"\n{res}")
                    # Intentar copiar fuente del primer run si existe
                    if el.runs and len(el.runs) > 1:
                        source_font = el.runs[0].font
                        run.font.name = source_font.name
                        run.font.size = source_font.size
                        run.font.bold = source_font.bold
                        run.font.italic = source_font.italic
                else: # Es una celda de tabla
                    first_para = el.paragraphs[0]
                    first_para.add_run(f"\n{res}")
                
                self.upd_prog()
            
            # Guardar progreso parcial
            doc.save(output_path)

    def translate_excel(self, orig, temp):
        wb = None
        try:
            wb = openpyxl.load_workbook(temp, data_only=True, keep_vba=True)
            new_wb = openpyxl.Workbook()

            batch_size = 15
            
            for idx, sn in enumerate(wb.sheetnames):
                if self.is_cancelled: break
                ws = wb[sn]
                nws = new_wb.active if idx == 0 else new_wb.create_sheet(sn)
                nws.title = sn

                for c, d in ws.column_dimensions.items(): nws.column_dimensions[c].width = d.width
                for r, d in ws.row_dimensions.items(): nws.row_dimensions[r].height = d.height

                max_row = ws.max_row
                max_col = ws.max_column

                pending_texts = []
                pending_cells = []
                
                for r in range(1, max_row + 1):
                    for c in range(1, max_col + 1):
                        if self.is_cancelled: return
                        cell, ncell = ws.cell(row=r, column=c), nws.cell(row=r, column=c)
                        
                        if not isinstance(cell, MergedCell) and self.is_valid_for_translation(cell.value):
                            pending_texts.append(str(cell.value))
                            pending_cells.append((cell.value, ncell))
                        else:
                            ncell.value = cell.value

                        if len(pending_texts) >= batch_size:
                            self.execute_translation_batch(pending_texts, pending_cells)
                            pending_texts, pending_cells = [], []

                        if cell.has_style:
                            ncell.font = copy.copy(cell.font)
                            ncell.border = copy.copy(cell.border)
                            ncell.fill = copy.copy(cell.fill)
                            ncell.alignment = Alignment(
                                    wrap_text=True, 
                                    horizontal=cell.alignment.horizontal or 'center', 
                                    vertical=cell.alignment.vertical or 'center'
                                )
                
                if pending_texts:
                    self.execute_translation_batch(pending_texts, pending_cells)

                for mr in ws.merged_cells.ranges: nws.merge_cells(str(mr))
                if hasattr(ws, '_images'):
                    for img in list(ws._images):
                        try:
                            new_img = copy.copy(img)
                            nws.add_image(new_img)
                        except Exception as e:
                            messagebox.showerror("Error", f"No se pudo copiar una imagen: {e}")

            target_lang_name = self.combo_to.get().upper()
            out_name = f"{target_lang_name}_BILINGUAL_{os.path.basename(orig)}"
            if orig.lower().endswith('.xls'): out_name += "x"
            new_wb.save(os.path.join(os.path.dirname(orig), out_name))
        except Exception as e:
            messagebox.showerror("Error Excel", str(e))
        finally:
            if wb: wb.close() # Importante para liberar memoria
            if "_TEMP_" in temp and os.path.exists(temp):
                try: os.remove(temp)
                except: pass
    
    def get_note(self):
        try:
            response = requests.get(f"{self.URL}/quote/loading", timeout=2)
            return response.json().get("text")
        except:
            return "Un documento menos, un beso más"
    
    def type_love_note(self):
        text = self.get_note()
        self.lbl_msg.configure(text="")
        def _type(i=0):
            if i <= len(text):
                self.lbl_msg.configure(text=text[:i])
                self.after(50, lambda: _type(i+1))
        _type()

    def upd_prog(self):
        self.done_units += 1
        total = self.total_units if self.total_units > 0 else 1
        p = self.done_units / total
        self.bar.set(p)

        if self.done_units % 50 == 0: 
            threading.Thread(target=self.type_love_note, daemon=True).start()

        elapsed = time.time() - self.start_time
        if self.done_units > 0 and self.total_units > 0:
            avg_time = elapsed / self.done_units
            rem = avg_time * (self.total_units - self.done_units)
            m, s = divmod(int(max(0, rem)), 60)
            self.lbl_eta.configure(text=f"ETA Global: {m:02d}:{s:02d}")
        self.update_idletasks()
    
    def show_love_note(self):
        if self.is_cancelled:
            return

        note = ctk.CTkToplevel(self)
        note.title("Hecho con cariño")
        note.geometry("350x250")
        note.attributes("-topmost", True)
        
                
        label = ctk.CTkLabel(
            note, 
            text=self.selected_text, 
            font=ctk.CTkFont(size=15, slant="italic"),
            justify="center"
        )
        label.pack(expand=True, padx=20, pady=20)
        
        btn = ctk.CTkButton(
            note, 
            text="Cerrar", 
            command=note.destroy, 
            fg_color="#db7093",
            hover_color="#c71585"
        )
        btn.pack(pady=(0, 20))

    def cleanup_janitor(self):
        processes = ["WINWORD.EXE", "EXCEL.EXE", "POWERPNT.EXE"]
        for proc in processes:
            try:
                subprocess.run(["taskkill", "/F", "/IM", proc, "/T"], 
                            stdout=subprocess.DEVNULL, 
                            stderr=subprocess.DEVNULL)
            except Exception:
                pass
        gc.collect()

    def safe_translate(self, text):
        if not text or not text.strip():
            return text
        
        clean_text = text.strip()
        if clean_text in self.translation_cache:
            return self.translation_cache[clean_text]
        
        try:
            translated = self.translate_with_retry(clean_text)
            self.translation_cache[clean_text] = translated
            self.api_counter += 1

            if self.api_counter % 10 == 0:
                self.save_cache()
            return translated
        except:
            return text
        
    def process_batch(self, text_list):
        if not text_list: return []

        final_translations = [None] * len(text_list)
        to_translate_api = []
        api_indices = []

        # 1. Filtro de Caché
        for i, text in enumerate(text_list):
            clean_text = self.normalize_text(text)
            if not clean_text:
                final_translations[i] = ""
                continue
            
            if clean_text in self.translation_cache:
                pure_translation = self.translation_cache[clean_text]
                
                # Aplicamos la lógica del checkbox al valor de la caché
                if self.keep_original.get():
                    final_translations[i] = f"{clean_text}\n{pure_translation}"
                else:
                    final_translations[i] = pure_translation
            else:
                to_translate_api.append(clean_text)
                api_indices.append(i)

        if not to_translate_api:
            return final_translations

        # 2. Delimitador
        delimiter = " [###] " 
        combined_text = delimiter.join(to_translate_api)
        
        try:
            translated_combined = self.translate_with_retry(combined_text)
            
            # 3. SPLIT con Regex
            api_results = re.split(r'\s?\[###\]\s?', translated_combined)
            
            # 4. Verificación de Integridad
            if len(api_results) != len(to_translate_api):
                api_results = [self.translate_with_retry(t) for t in to_translate_api]

            # 5. Reinsertar y alimentar caché (LÓGICA CORREGIDA)
            for i, res in enumerate(api_results):
                if i < len(api_indices):
                    target_idx = api_indices[i]
                    pure_translation = res.strip() # Traducción limpia de la API
                    orig_text = to_translate_api[i]
                    
                    # SIEMPRE guardamos la traducción PURA en la caché
                    # Esto hace que la caché sea independiente del formato visual
                    self.translation_cache[orig_text] = pure_translation
                    
                    # Aplicamos formato para el documento actual
                    if self.keep_original.get():
                        final_translations[target_idx] = f"{orig_text}\n{pure_translation}"
                    else:
                        final_translations[target_idx] = pure_translation
            
            self.api_counter += 1
            
        except Exception as e:
            messagebox.showerror("Error", f"Error API Batch: {e}")
            for idx in api_indices:
                if final_translations[idx] is None:
                    final_translations[idx] = text_list[idx]

        return final_translations
        
    def execute_translation_batch(self, texts, cell_objects):
        translations = self.process_batch(texts)

        for i, (original_val, ncell) in enumerate(cell_objects):
            try:
                # Usamos el resultado de process_batch directamente
                # porque esa función ya aplicó la lógica del checkbox.
                final_val = translations[i] if i < len(translations) else ""
                
                # Asignación limpia al Excel
                ncell.value = final_val if final_val is not None else original_val
                
            except Exception as e:
                messagebox.showerror("Error", f"Error en celda {ncell.coordinate}: {e}")
                ncell.value = original_val
                
            self.upd_prog()
        
        self.save_cache()
    
    def is_valid_for_translation(self, value):
        # Si es None o un string vacío/solo espacios, fuera.
        if value is None:
            return False
            
        val_str = str(value).strip()
        if not val_str:
            return False
            
        # Ignorar si es solo un número (ya sea int o float)
        try:
            float(val_str)
            return False
        except ValueError:
            pass

        # Solo si contiene al menos una letra (evita celdas que son solo "!!!" o "---")
        return any(c.isalpha() for c in val_str)
    
    def translate_with_retry(self, text, retries=10):
        for i in range(retries):
            try:
                # Tu lógica de traducción aquí
                return self.translator.translate(text)
            except Exception:
                sleep(2 ** i)

    def get_cache_filename(self):
        return f"memory_{self.current_src_lang}_{self.current_tgt_lang}.json"

    def load_cache(self):
        fname = self.get_cache_filename()
        if os.path.exists(fname):
            try:
                with open(fname, "r", encoding="utf-8") as f:
                    return json.load(f)
            except: return {}
        return {}

    def save_cache(self):
        """Guarda la memoria segmentada y limpia."""
        fname = self.get_cache_filename()
        
        # Limpiamos la memoria antes de guardar
        clean_memory = self.deduplicate_memory(self.translation_cache)
        
        try:
            with open(fname, "w", encoding="utf-8") as f:
                json.dump(clean_memory, f, ensure_ascii=False, indent=4)
            # Actualizamos la memoria en vivo con la versión limpia
            self.translation_cache = clean_memory
        except Exception as e:
            messagebox.showerror("Error",f"Error guardando caché {fname}: {e}")

    def normalize_text(self, text):
        """Unifica espacios y limpia el texto para evitar duplicados."""
        if text is None: return ""
        # Convierte a string, quita espacios raros y unifica a un solo espacio entre palabras
        return re.sub(r'\s+', ' ', str(text)).strip()
    
    def deduplicate_memory(self, memory_dict):
        """Limpia claves duplicadas por espacios o normalización antes de guardar."""
        new_memory = {}
        for key, value in memory_dict.items():
            # Normalizamos la clave para unificar "Día " con "Día"
            norm_key = self.normalize_text(key)
            if norm_key and norm_key not in new_memory:
                new_memory[norm_key] = value
        return new_memory
    
    def translate_shapes_via_win32(self, file_path):
        import pythoncom
        pythoncom.CoInitialize()
        app = None
        try:
            app = win32.Dispatch('Excel.Application')
            app.Visible = False
            app.DisplayAlerts = False
            wb = app.Workbooks.Open(os.path.abspath(file_path))

            for sheet in wb.Sheets:
                for shape in sheet.Shapes:
                    try:
                        if shape.HasTextFrame:
                            original_text = shape.TextFrame2.TextRange.Text
                            clean_text = self.normalize_text(original_text)
                            
                            if self.is_valid_for_translation(clean_text):
                                # Obtenemos la traducción (ya normalizada por process_batch)
                                translated_text = self.safe_translate(clean_text)
                                
                                # --- LÓGICA DE MANTENER ORIGINAL ---
                                # Si safe_translate ya trae el original por estar en caché, 
                                # esta lógica se mantiene consistente.
                                if self.keep_original.get():
                                    # Verificamos si el texto ya contiene el original para no duplicarlo
                                    if original_text not in translated_text:
                                        shape.TextFrame2.TextRange.Text = f"{original_text}\n{translated_text}"
                                    else:
                                        shape.TextFrame2.TextRange.Text = translated_text
                                else:
                                    shape.TextFrame2.TextRange.Text = translated_text
                                
                                shape.TextFrame.AutoSize = True
                    except Exception:
                        continue

            wb.Save() # Usamos Save() para sobreescribir el temporal
            wb.Close()
        except Exception as e:
            messagebox.showerror("Error", f"Error procesando Shapes: {e}")
        finally:
            if app:
                app.Quit()
            pythoncom.CoUninitialize()

    def translate_pptx_shapes_via_win32(self, file_path):
        """Translates PPT shapes preserving exact styles and fixing Read-Only errors."""
        import pythoncom
        import stat
        pythoncom.CoInitialize()
        app = None
        try:
            # 1. Clear Read-Only attribute from the file immediately
            abs_path = os.path.abspath(file_path)
            if os.path.exists(abs_path):
                os.chmod(abs_path, stat.S_IWRITE)

            app = win32.Dispatch('PowerPoint.Application')
            # Open without window to avoid focus issues
            pres = app.Presentations.Open(abs_path, WithWindow=False)

            # Define the final output path to avoid saving over the read-only source
            target_lang_name = self.combo_to.get().upper()
            out_name = f"{target_lang_name}_BILINGUAL_{os.path.basename(file_path).replace('TEMP_', '')}"
            output_path = os.path.join(os.path.dirname(file_path), out_name)
            abs_output_path = os.path.abspath(output_path)

            for slide in pres.Slides:
                for shape in slide.Shapes:
                    self.process_win32_shape_recursive(shape)

            # 2. Use SaveAs to the final destination instead of saving the temp file
            # FileFormat 1 is a standard PPTX
            pres.SaveAs(abs_output_path) 
            pres.Close()

        except Exception as e:
            messagebox.showerror("Error", f"Critical Win32 PPT Error: {e}")
            # If SaveAs fails, try one last resort: Save to user's Documents
            try:
                emergency_path = os.path.join(os.path.expanduser("~"), "Documents", "Recovered_Translation.pptx")
                pres.SaveAs(emergency_path)
            except: pass
        finally:
            if app: 
                app.Quit()
            pythoncom.CoUninitialize()

    def process_win32_shape_recursive(self, shape):
        """Recursively dive into grouped items to find text and apply bilingual logic."""
        # 3. Handle Grouped Shapes (msoGroup = 6)
        if shape.Type == 6: 
            try:
                for item in shape.GroupItems:
                    self.process_win32_shape_recursive(item)
            except: pass
            return

        try:
            if shape.HasTextFrame and shape.TextFrame.HasText:
                orig_range = shape.TextFrame.TextRange
                original_text = orig_range.Text
                clean_text = self.normalize_text(original_text)
                
                if self.is_valid_for_translation(clean_text):
                    # Get translation (uses cache if available)
                    translated_text = self.safe_translate(clean_text)
                    
                    if self.keep_original.get():
                        # .InsertAfter(text) in Win32 inherits the style of the last character
                        # which is perfect for maintaining font, size, and color.
                        new_range = orig_range.InsertAfter(f"\r{translated_text}")
                        
                        # Optional: If you want the translation to be slightly smaller or gray:
                        # new_range.Font.Size = max(orig_range.Font.Size - 2, 8)
                        # new_range.Font.Color.RGB = 0x808080 
                    else:
                        # Replace but keep original styling from the first character
                        orig_range.Text = translated_text
                    
                    # Ensure the text box grows to fit the new text
                    shape.TextFrame.AutoSize = 1 # ppAutoSizeShapeToFitText
        except Exception as e:
            messagebox.showerror("Error", f"Error in shape logic: {e}")
    
    def get_all_shapes(self, shapes):
        """Recursively finds all shapes, even those nested in groups."""
        all_shps = []
        for shape in shapes:
            if shape.shape_type == 6:  # 6 is the constant for a GroupShape
                all_shps.extend(self.get_all_shapes(shape.shapes))
            else:
                all_shps.append(shape)
        return all_shps

    def remove_file_from_queue(self, path_to_remove, frame_to_destroy):
        """Removes a file from the internal logic and destroys its UI row."""
        # 1. Remove from the path set (allows re-adding later)
        if path_to_remove in self.files_path_set:
            self.files_path_set.remove(path_to_remove)
        
        # 2. Remove from the processing list
        self.files_to_process = [f for f in self.files_to_process if f["path"] != path_to_remove]
        
        # 3. Destroy the UI frame
        frame_to_destroy.destroy()
        
        # 4. Disable 'Go' button if queue is empty
        if not self.files_to_process:
            self.btn_go.configure(state="disabled")

if __name__ == "__main__":
    app = DocTranslatorPro()
    app.mainloop()